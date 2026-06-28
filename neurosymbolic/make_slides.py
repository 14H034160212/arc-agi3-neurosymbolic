#!/usr/bin/env python3
"""Generate the project deck (ARC-AGI-3 neuro-symbolic agent) as an editable .pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x12, 0x2A, 0x4A); BLUE = RGBColor(0x1F, 0x6F, 0xB2)
GREEN = RGBColor(0x1B, 0x7A, 0x43); GREY = RGBColor(0x44, 0x4A, 0x55)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Microsoft YaHei"
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def _set(tf_para, size, color=GREY, bold=False):
    tf_para.font.size = Pt(size); tf_para.font.bold = bold
    tf_para.font.color.rgb = color; tf_para.font.name = FONT


def band(slide, color, top, height):
    s = slide.shapes.add_shape(1, 0, top, W, height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def title_slide(title, subtitle, tag):
    s = prs.slides.add_slide(BLANK)
    band(s, NAVY, 0, H)
    box = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.2)).text_frame
    box.word_wrap = True
    p = box.paragraphs[0]; p.text = title; _set(p, 40, WHITE, True)
    p2 = box.add_paragraph(); p2.text = subtitle; _set(p2, 22, RGBColor(0xCF, 0xE0, 0xF5)); p2.space_before = Pt(14)
    p3 = box.add_paragraph(); p3.text = tag; _set(p3, 15, RGBColor(0x9F, 0xC2, 0xE8)); p3.space_before = Pt(22)


def content(title, bullets, foot=None):
    s = prs.slides.add_slide(BLANK)
    band(s, NAVY, 0, Inches(1.1))
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(12.1), Inches(0.8)).text_frame
    t.word_wrap = True; p = t.paragraphs[0]; p.text = title; _set(p, 26, WHITE, True)
    body = s.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(12.0), Inches(5.6)).text_frame
    body.word_wrap = True
    for i, b in enumerate(bullets):
        lvl = 0; txt = b; col = GREY; sz = 18; bold = False
        if isinstance(b, tuple):
            txt, lvl = b[0], b[1]
            if len(b) > 2: col = b[2]
            if len(b) > 3: bold = b[3]
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = ("• " if lvl == 0 else "– ") + txt
        _set(p, sz - (2 if lvl else 0), col if col != GREY else (NAVY if lvl == 0 else GREY), bold or lvl == 0 and False)
        p.level = lvl; p.space_after = Pt(7)
        if lvl == 0: p.font.color.rgb = NAVY if bold else GREY; p.font.bold = bold
    if foot:
        f = s.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(12), Inches(0.4)).text_frame
        p = f.paragraphs[0]; p.text = foot; _set(p, 12, BLUE, True)
    return s


def table_slide(title, headers, rows, foot=None, col_w=None):
    s = prs.slides.add_slide(BLANK)
    band(s, NAVY, 0, Inches(1.1))
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(12.1), Inches(0.8)).text_frame
    p = t.paragraphs[0]; p.text = title; _set(p, 26, WHITE, True)
    nr, nc = len(rows) + 1, len(headers)
    gt = s.shapes.add_table(nr, nc, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.5 * nr)).table
    if col_w:
        for i, w in enumerate(col_w): gt.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = gt.cell(0, j); c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = BLUE
        pr = c.text_frame.paragraphs[0]; _set(pr, 15, WHITE, True)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i + 1, j); c.text = val
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT
            pr = c.text_frame.paragraphs[0]
            col = GREEN if ("✅" in val or "7/7" in val) else GREY
            _set(pr, 13, col, "✅" in val or j == 0)
    if foot:
        f = s.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(12), Inches(0.4)).text_frame
        p = f.paragraphs[0]; p.text = foot; _set(p, 12, BLUE, True)
    return s


# ───────────────────────────── slides ─────────────────────────────
title_slide("ARC-AGI-3 神经符号 Agent",
            "从复现 → 纯像素源码无关 7/7 → 跨游戏 / 跨模态泛化",
            "本地零成本可复现 · 代码 + 报告已开源 · 2026")

content("1 · 背景：ARC-AGI-3 考什么", [
    ("64×64 彩色网格的小游戏：规则隐藏、没有说明书", 0),
    ("agent 只能靠动作交互，自己发现规则并逐关通关", 0),
    ("考的是「在陌生环境里临场获取技能」的流体智能——纯记忆/堆数据难奏效", 0),
    ("现状参照：Symbolica 用 harness（前沿模型 Opus + 脚手架），公开集 36.08%", 0),
    ("本质 = 别人训好的大模型 + 脚手架；智能在模型", 1),
], foot="问题：能不能不靠大模型的「黑箱智能」，让 agent 自己把隐藏规则推出来？")

content("2 · 三条路线 & 我们的选择", [
    ("① 自进化发现 workflow —— 搜索空间大、易过拟合", 0),
    ("② 把 workflow 训进模型（端到端）—— 天花板高但烧算力、撞泛化墙", 0),
    ("③ 模型写代码（PDDL/Prolog）做符号推理 ← 选它", 0, NAVY, True),
    ("代码数据多、泛化好、搜索可卸载给求解器（有正确性保证）", 1),
    ("核心论点：让 LLM 只「把规则写成可执行代码」，搜索交给符号规划器 → 连免费/小模型都能用", 0),
    ("借鉴 ABPR（arXiv:2603.20334）：程序即假设 + 算法化调试，首次落到交互式 ARC-AGI-3", 1),
])

content("3 · 方法总览：感知 + 交互 + 规划 + 验证", [
    ("读状态 → 把动力学写成代码 → 符号规划器搜索 → 在真环境自动验证", 0, NAVY, True),
    ("↑ APD 定位错误子句 + 定向修正（验证信号 = 规划解能否真通关）", 1),
    ("不需要训练即可起步；训练（蒸馏归纳轨迹）是可选第二阶段", 0),
    ("全程本地、零成本：引擎离线 + 免费本地 gpt-oss-120b", 0),
    ("逆向工程发现 ls20 = 配送拼图：携带物体(形状/颜色/旋转)→ 改造站调属性 → 送到匹配槽", 0),
    ("隐藏状态：能量(每步−1)、命数(3)；确定性；像素是非平凡渲染（感知是真正难点）", 1),
])

table_slide("4 · 结果一 & 二：符号规划 7/7 + 免费 LLM 归纳 3/7", [
    "成果", "结果", "说明"], [
    ["符号规划（引擎态）", "✅ 7/7", "验证过的动力学模型 + 子目标规划器，7 关全解（最优）"],
    ["对比基线", "—", "规划器 14 步最优 / Opus 探索 42 步 / 本地裸 LLM 0 通关"],
    ["免费 LLM 归纳", "✅ 3/7", "本地 gpt-oss 从观测自主归纳动力学代码 + 规划式 APD 反馈"],
    ["关键洞见", "—", "规划 trivial；难点在「归纳 + 像素感知」，不在搜索"],
], col_w=[3.2, 1.6, 7.1])

content("5 · 结果三（核心）：纯像素源码无关 Agent —— 7/7", [
    ("只用渲染画面 + 通关/失败反馈，端到端解出全部 7 关（逐关串联）", 0, GREEN, True),
    ("与「读引擎状态的符号规划器 7/7」持平，但只用像素！", 0, NAVY, True),
    ("难度逐关攀升，全部自己学会：", 0),
    ("L0 旋转×1 · L1 旋转×3 · L2 颜色+旋转 · L3 · L4 形状+颜色+旋转", 1),
    ("L5 两个堆叠的槽（多槽，先解上槽打通下槽）· L6 形状+颜色+旋转（胜利=WIN 态）", 1),
    ("全部 7 关 ~3 分钟跑完；关卡靠「赢了才进下一关」串联（RESET 回到 L0）", 0),
], foot="uv run python neurosymbolic/source_free_agent.py  → 一键复现 7/7")

content("6 · 一路被真实 bug 逼出来的关键机制", [
    ("颜色重载 → 按颜色识别不可行（色9 既是玩家调色板又是槽标记）", 0),
    ("航位推算定位玩家：用「已学动作方向」而非固定颜色质心 → 对重绘/变色免疫", 0),
    ("方向用排除法学一次（避免起点被墙挡的「401 幽灵格」bug）", 0),
    ("站点从「固定 HUD 仓库面板」读取：配置画在角落而非玩家身上", 0),
    ("能量是隐藏变量：64×64 画面根本不渲染 → 用「狂走到死=被瞬移回起点，步数即能量」反推", 0, BLUE, True),
    ("补给点可见(色11) → 能量感知 Dijkstra + 提前补给；死亡会清空配置 → 活着也是计划的一部分", 1),
    ("多槽：基于「投放是否真把玩家移进被挡格」判定槽解（对遮挡免疫）+ 标记轮询优先级", 0),
])

content("7 · 方法论价值：这正是 ARC-AGI-3 想考的能力", [
    ("agent 只凭画面 + 反馈，从零自己发现了：", 0, NAVY, True),
    ("位置、动作基(上下左右)、站点、目标槽、多步配置(形状/颜色/旋转)", 1),
    ("一个完全不可见的隐藏资源(能量) + 补给点 + 多槽求解顺序", 1, BLUE, False),
    ("= 「在陌生环境里从后果反推隐藏规则」的流体智能", 0, GREEN, True),
    ("验证信号应是「规划解能否真通关」，不是逐转移准确率", 0),
    ("每种机制至少要观测一次才能归纳 → 指向主动探索", 0),
])

content("8 · 泛化：把方法论抽成 game-agnostic 骨架", [
    ("把单游戏解法切成两半（source_free_core.py）：", 0, NAVY, True),
    ("通用算子（可迁移，只依赖 Env=reset/step/render/win-lose/score）：", 0, BLUE, True),
    ("探测学动作语义 · 靠运动找/追踪主体 · reset+回放建图 · 隐藏资源预言机", 1),
    ("靠交互而非外观发现物体角色(transformer/gate/replenisher) · 胜负验证的搜索", 1),
    ("每游戏适配器（小、需归纳）：谁是 agent / 怎么读目标 / 哪些格可交互", 0, BLUE, True),
    ("通用求解器自己发现 transformer，并把 ls20 L0 用 16 步解出（零 ls20 知识）", 0, GREEN, True),
], foot="ls20 7/7 = 存在性证明：只要每游戏抽象对了，源码无关符号求解就是完整的")

table_slide("9 · 跨游戏 + 双模态：同一套循环解 3 个游戏的 L0", [
    "游戏", "自动识别模态", "通用算子结果（无每游戏代码）"], [
    ["ls20", "movement 移动", "✅ 通用求解器靠发现解 L0(16 步)；完整版 7/7"],
    ["ft09", "click 点击", "✅ L0 用 4 次点击解出"],
    ["vc33", "click 点击", "✅ L0 用 3 次点击解出"],
], foot="detect_modality 自动选移动 vs 点击算子族 · uv run python neurosymbolic/source_free_core.py",
    col_w=[1.8, 3.2, 6.9])

content("10 · LLM 适配器：免费本地大模型做「每游戏归纳」", [
    ("当自动发现不够时，LLM 看几帧 + 探测效果 → 输出该游戏的适配器", 0),
    ("已接通本地 gpt-oss-120b（免费，Ollama）并真实验证：", 0, NAVY, True),
    ("给它「方向动作无反应 + 这些格点击有反应」→ 它答 modality=click 并给出正确理由", 1, GREEN, True),
    ("call_llm 注入式、provider 无关：换任意 OpenAI 兼容 API 即可", 0),
    ("这正是神经符号方案的落点：LLM 负责「在陌生游戏把抽象 induce 出来」，", 0),
    ("通用循环负责「规划 + 验证」（有正确性保证）", 1),
])

content("11 · 整体版图 & 下一步", [
    ("① 符号规划 7/7（证明 ls20 可解）", 0, GREEN, True),
    ("② 免费 LLM 归纳 3/7（证明 LLM 能写动力学代码）", 0, GREEN, True),
    ("③ 纯像素源码无关 agent 7/7（证明只凭画面+反馈能自发现全部隐藏规则）", 0, GREEN, True),
    ("④ game-agnostic 骨架 + 双模态（3 游戏 L0 全解）+ LLM 适配器接通", 0, GREEN, True),
    ("下一步：", 0, NAVY, True),
    ("更聪明的搜索/结构去攻深层关卡（vc33/ft09 L1+）", 1),
    ("LLM 适配器在更多 ARC-AGI-3 公开游戏上量跨游戏迁移率（对标 Symbolica 36%）", 1),
])

content("12 · 结论", [
    ("用真实实验证明：神经符号路线在 ARC-AGI-3 上成立，且可迁移", 0, NAVY, True),
    ("纯像素、只靠画面+反馈 → 自己发现位置/站点/槽/多步配置/隐藏资源 → ls20 7/7", 0),
    ("方法论抽成通用算子 + 小适配器 → 跨游戏、跨模态（移动+点击）解出 3 个游戏 L0", 0),
    ("免费本地 LLM 已接通，做「每游戏归纳」那一环", 0),
    ("全部可复现：代码 + 报告 + 本 PPT 已开源", 0, GREEN, True),
], foot="arc-neurosymbolic-repo · neurosymbolic/{source_free_agent.py, source_free_core.py, REPORT.md}")

import sys
out = sys.argv[1] if len(sys.argv) > 1 else "ARC_AGI3_neurosymbolic.pptx"
prs.save(out)
print(f"saved {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
